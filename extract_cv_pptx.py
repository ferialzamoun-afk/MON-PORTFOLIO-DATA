from pptx import Presentation

# Charger la présentation
prs = Presentation('Ferial-Zamoun.pptx')

# Extraire tout le texte
def extract_text(prs):
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Écrire dans un fichier UTF-8
with open("cv_extrait.txt", "w", encoding="utf-8") as f:
    f.write(extract_text(prs))

print("Extraction terminée. Voir cv_extrait.txt.")
